"""
data_loader.py — Multi-format file text extractor.

Supported formats:
  PDF         → PyMuPDF (pymupdf)
  Word        → python-docx (.docx)
  Excel       → openpyxl (.xlsx) / xlrd (.xls)
  PowerPoint  → python-pptx (.pptx)
  Images      → Pillow + pytesseract OCR (.png, .jpg, .jpeg, .bmp, .tiff, .webp)
  Text        → plain read (.txt, .csv, .md)

Install requirements:
  pip install pymupdf python-docx openpyxl python-pptx pillow pytesseract
  Also install Tesseract OCR engine: https://github.com/UB-Mannheim/tesseract/wiki
"""

import os

# ═══════════════════════════════════════════
# MAIN DISPATCHER
# ═══════════════════════════════════════════

def load_file(file_path: str) -> str:
    """
    Auto-detect file type by extension and extract text.
    Returns extracted text string, or empty string on failure.
    """
    ext = os.path.splitext(file_path)[1].lower()

    loaders = {
        ".pdf":  _load_pdf,
        ".docx": _load_docx,
        ".doc":  _load_docx,
        ".xlsx": _load_excel,
        ".xls":  _load_excel,
        ".pptx": _load_pptx,
        ".ppt":  _load_pptx,
        ".txt":  _load_text,
        ".md":   _load_text,
        ".csv":  _load_text,
        ".png":  _load_image,
        ".jpg":  _load_image,
        ".jpeg": _load_image,
        ".bmp":  _load_image,
        ".tiff": _load_image,
        ".tif":  _load_image,
        ".webp": _load_image,
    }

    loader = loaders.get(ext)
    if not loader:
        print(f"[WARNING] Unsupported file type: {ext}")
        return ""

    try:
        text = loader(file_path)
        print(f"[LOADER] {ext} → {len(text.split())} words extracted from {os.path.basename(file_path)}")
        return text
    except Exception as e:
        print(f"[ERROR] Failed to load {file_path}: {e}")
        return ""


# ═══════════════════════════════════════════
# PDF  (PyMuPDF)
# ═══════════════════════════════════════════

def _load_pdf(file_path: str) -> str:
    import pymupdf as fitz
    doc = fitz.open(file_path)
    pages = []
    for page_num, page in enumerate(doc):
        try:
            t = page.get_text()
            if t and t.strip():
                pages.append(t)
        except Exception as e:
            print(f"[WARNING] PDF page {page_num} skipped: {e}")
    doc.close()
    return "\n".join(pages)


# Backward-compatible alias used by old code
def load_pdf(file_path: str) -> str:
    return _load_pdf(file_path)


# ═══════════════════════════════════════════
# WORD  (python-docx)
# ═══════════════════════════════════════════

def _load_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    parts = []

    # Paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)

    return "\n".join(parts)


# ═══════════════════════════════════════════
# EXCEL  (openpyxl)
# ═══════════════════════════════════════════

def _load_excel(file_path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []

    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None and str(cell).strip())
            if row_text:
                parts.append(row_text)

    wb.close()
    return "\n".join(parts)


# ═══════════════════════════════════════════
# POWERPOINT  (python-pptx)
# ═══════════════════════════════════════════

def _load_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {slide_num}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            # Tables inside PPT
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)

    return "\n".join(parts)


# ═══════════════════════════════════════════
# PLAIN TEXT  (.txt, .csv, .md)
# ═══════════════════════════════════════════

def _load_text(file_path: str) -> str:
    encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            continue
    return ""


# ═══════════════════════════════════════════
# IMAGES  (Pillow + Tesseract OCR)
# ═══════════════════════════════════════════

def _load_image(file_path: str) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text
    except ImportError:
        print("[WARNING] pytesseract/Pillow not installed. Image OCR skipped.")
        return ""
    except Exception as e:
        print(f"[WARNING] OCR failed for {file_path}: {e}")
        return ""